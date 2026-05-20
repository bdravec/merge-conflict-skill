"""Insert a baseline pass/fail slide into the 2026-05-22 review deck (closes #61).

Inserts a new slide at position 8 (0-indexed: 7) anchored on
`docs/figures/baseline_pass_fail.png`, with headline numbers in a right-side
text box. Existing slides shift down by one.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPO = Path(__file__).resolve().parent.parent
DECK = Path("/home/baebs/thesis/review_mtgs/Thesis_Review_2026-05-22.pptx")
FIGURE = REPO / "docs/figures/baseline_pass_fail.png"
INSERT_AT = 7  # 0-indexed; slide will become slide 8 in 1-indexed PPT view

TITLE_A = "ConGra baseline — python-tiny no-skill"
TITLE_B = "Pass/fail tiers (Zhang et al. 2024 thresholds)"
TOC_HEADER = "3. Methodology"

BULLETS = [
    ("Qwen3-8B", "29.2% solved · 6.5% failed", True),
    ("Apertus-8B", "21.5% solved · 7.5% failed", True),
    ("n = 3,597", "7 complexity buckets", False),
    ("solved = max(edit, winn) > 0.8", "failed = ≤ 0.05", False),
    ("Hardest bucket bimodal", "text+sytx+func: peaks at 0 and 1.0", False),
]


def find_layout(prs, name):
    for sm in prs.slide_masters:
        for lay in sm.slide_layouts:
            if lay.name == name:
                return lay
    raise KeyError(name)


def set_placeholder_text(slide, idx, text, font_size=None, bold=False, color=None):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            if font_size is not None:
                run.font.size = Pt(font_size)
            if bold:
                run.font.bold = True
            if color is not None:
                run.font.color.rgb = color
            return
    raise KeyError(f"placeholder idx={idx} not found")


def build():
    assert DECK.exists(), DECK
    assert FIGURE.exists(), FIGURE

    prs = Presentation(str(DECK))
    layout = find_layout(prs, "5b: Text-Folie 2 Spalten")
    slide = prs.slides.add_slide(layout)

    # Title (red) and subtitle (black) using layout's placeholder schema.
    set_placeholder_text(slide, 20, TOC_HEADER)
    set_placeholder_text(slide, 0, TITLE_A)
    set_placeholder_text(slide, 11, TITLE_B)

    # Remove the two body content placeholders (left + right) — we'll place
    # the image and a custom text box ourselves so the chart can be wider
    # than the default 4in column.
    for idx_to_drop in (1, 21):
        for ph in list(slide.placeholders):
            if ph.placeholder_format.idx == idx_to_drop:
                sp = ph._element
                sp.getparent().remove(sp)
                break

    # Image: 5.6in wide centered on the left, aspect preserved from PNG (≈2.5:1).
    img_left = Emu(540000)
    img_top = Emu(1650000)
    img_width = Emu(5200000)
    slide.shapes.add_picture(
        str(FIGURE), img_left, img_top, width=img_width
    )

    # Right text box.
    txt_left = Emu(5900000)
    txt_top = Emu(1700000)
    txt_width = Emu(3000000)
    txt_height = Emu(3000000)
    tb = slide.shapes.add_textbox(txt_left, txt_top, txt_width, txt_height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for head, sub, bold_head in BULLETS:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
        r1 = p.add_run()
        r1.text = head
        r1.font.size = Pt(14)
        r1.font.bold = bold_head
        if bold_head:
            r1.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        r2 = p.add_run()
        r2.text = "\n" + sub
        r2.font.size = Pt(12)
        r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Move the just-appended slide from the end to position INSERT_AT.
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    moving = slides_list[-1]
    xml_slides.remove(moving)
    xml_slides.insert(INSERT_AT, moving)

    out = Path("/tmp/Thesis_Review_2026-05-22.pptx.new")
    prs.save(str(out))

    # Sanity check.
    verify = Presentation(str(out))
    assert len(verify.slides) == 29, f"expected 29 slides, got {len(verify.slides)}"
    titles = []
    for s in verify.slides:
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                titles.append(sh.text_frame.text.strip().split("\n")[0][:60])
                break
    print(f"slide count: {len(verify.slides)}")
    for i, t in enumerate(titles, 1):
        marker = " <-- NEW" if i == INSERT_AT + 1 else ""
        print(f"  {i:>2}. {t}{marker}")
    print(f"wrote {out}")


if __name__ == "__main__":
    build()
