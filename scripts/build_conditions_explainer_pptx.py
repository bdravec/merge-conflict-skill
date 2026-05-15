"""
build_conditions_explainer_pptx.py — 2-slide explainer for pilot conditions (issue #52).

Source content: docs/conditions_explained.md.

Layout `5a: Text-Folie 2 Spalten` placeholder indices:
  idx=20  TOC header
  idx=0   Titel A (red)
  idx=11  Titel B (black subtitle)
  idx=16  Left column heading
  idx=17  Left column body
  idx=18  Right column heading
  idx=19  Right column body

Slide 2 drops the 4 column placeholders and adds a python-pptx table for
the 3 conditions x 2 slots mapping, plus a text box below for invariants
and the naming convention note.

python-pptx slide-deletion gotcha: removing from sldIdLst alone leaves
orphan parts; we drop the relationship too (delete_all_slides).
"""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn


TEMPLATE = Path("/home/baebs/thesis/review_mtgs/PPTVorlage-UniBe_ger.pptx")
DEFAULT_OUT = Path("/home/baebs/thesis/review_mtgs/conditions_explainer.pptx")


def find_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(f"Layout {name!r} not found in any slide master.")


def set_ph_text(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text_frame.text = text
            return
    raise KeyError(f"No placeholder with idx={idx} on slide.")


def set_ph_lines(slide, idx, lines):
    """Write multiple paragraphs into a placeholder, one per line."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.text = lines[0]
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line
            return
    raise KeyError(f"No placeholder with idx={idx} on slide.")


def drop_placeholder(slide, idx):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == idx:
            sp = ph._element
            sp.getparent().remove(sp)
            return


def delete_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    pres_part = prs.part
    for sld_id in list(sldIdLst):
        rId = sld_id.get(qn("r:id"))
        sldIdLst.remove(sld_id)
        pres_part.rels.pop(rId)


def add_slide_1_architecture(prs):
    layout = find_layout(prs, "5a: Text-Folie 2 Spalten")
    slide = prs.slides.add_slide(layout)
    set_ph_text(slide, 20, "Pilot conditions reference")
    set_ph_text(slide, 0, "Chat-message architecture")
    set_ph_text(slide, 11, "Each LLM call sends two messages: a system + a user message")

    set_ph_text(slide, 16, "system slot")
    set_ph_lines(slide, 17, [
        "Purpose: persona, rules, constraints.",
        "Sent before any conversation turn.",
        "Model treats it as authoritative setup —",
        "post-training conditions the model to weight",
        "system content more heavily than user content.",
        "",
        "ConGra default:",
        "\"You are an expert in code merge conflicts,",
        "providing the merged code based on the",
        "conflict and its context.\"",
    ])

    set_ph_text(slide, 18, "user slot")
    set_ph_lines(slide, 19, [
        "Purpose: the task or question from the human.",
        "Model treats it as the request to act on.",
        "",
        "ConGra CoT template:",
        "• Numbered reasoning steps",
        "  (1. understand cause → 2. decide → 3. emit",
        "   merged code in fenced ``` block)",
        "• Conflict payload",
        "  (language hint + conflict_context block +",
        "   conflict_text block)",
        "",
        "CoT = chain-of-thought: instruct the model to",
        "reason step-by-step before answering.",
    ])
    return slide


def add_slide_2_conditions(prs):
    layout = find_layout(prs, "5a: Text-Folie 2 Spalten")
    slide = prs.slides.add_slide(layout)
    set_ph_text(slide, 20, "Pilot conditions reference")
    set_ph_text(slide, 0, "The three pilot conditions")
    set_ph_text(slide, 11, "Each condition = a different configuration of the two slots")

    for idx in (16, 17, 18, 19):
        drop_placeholder(slide, idx)

    rows, cols = 4, 3
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(12.3)
    height = Inches(2.4)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(4.6)
    table.columns[2].width = Inches(5.5)

    headers = ["Condition", "system slot", "user slot"]
    cells = [
        ["no-skill",
         "ConGra default sysprompt\n(\"You are an expert in code merge conflicts…\")",
         "ConGra CoT template"],
        ["skill-vX-sys",
         "SKILL.md vX body\n(ConGra default is dropped — overwritten)",
         "ConGra CoT template"],
        ["skill-vX-user",
         "ConGra default sysprompt\n(unchanged)",
         "SKILL.md vX body + ConGra CoT template\n(skill prepended)"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)
                r.font.bold = True
    for ri, row in enumerate(cells, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)

    txbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.0),
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    tf.text = "Invariants"
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(13)

    invariants = [
        "• The system slot is never empty — it holds either the ConGra default or the SKILL.md body.",
        "• The user slot always carries the ConGra CoT template (instructions + payload).",
        "• The skill appears in exactly one slot per condition: nowhere, or in -sys, or in -user.",
    ]
    for line in invariants:
        p = tf.add_paragraph()
        p.text = line
        for r in p.runs:
            r.font.size = Pt(11)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Naming"
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(13)

    p = tf.add_paragraph()
    p.text = (
        "The -sys / -user suffix names the slot SKILL.md was injected into, "
        "not which slot is empty. Even in skill-vX-user, the system slot still "
        "contains ConGra's default sysprompt."
    )
    for r in p.runs:
        r.font.size = Pt(11)

    return slide


def main():
    ap = argparse.ArgumentParser(
        description="Build the 2-slide pilot-conditions explainer pptx.",
    )
    ap.add_argument("--template", type=Path, default=TEMPLATE,
                    help="Uni Bern template pptx path")
    ap.add_argument("--output", type=Path, default=DEFAULT_OUT,
                    help="Output pptx path")
    args = ap.parse_args()

    if not args.template.is_file():
        ap.error(f"Template not found: {args.template}")

    prs = Presentation(str(args.template))
    n_orig = len(prs.slides)
    print(f"Loaded template: {args.template} ({n_orig} slides)")

    delete_all_slides(prs)
    print("  -> cleared all template slides")

    add_slide_1_architecture(prs)
    print("  -> added slide 1: Chat-message architecture")

    add_slide_2_conditions(prs)
    print("  -> added slide 2: The three pilot conditions")

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output))
    print(f"Wrote: {args.output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
