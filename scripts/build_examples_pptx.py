"""
build_examples_pptx.py — populate conflict_1-6_examples.pptx for the
2026-05-22 review appendix (issue #49).

Deck (8 slides):
  1a  Conflict 1 — Apertus      [split, 5-col v1/v2/v2.1 comparison]
  1b  Conflict 1 — Qwen3        [split]
  2   Conflict 2                [4-col Apertus/Qwen3 side-by-side]
  3   Conflict 3                [4-col]
  4a  Conflict 4 — Apertus      [split]
  4b  Conflict 4 — Qwen3        [split]
  5   Conflict 5                [4-col]
  6   Conflict 6                [4-col]

Typography (Barbara 2026-05-14):
  Title (red):      Arial 28pt — "Conflict N" or "Conflict N — <Model>"
  Subtitle (black): Arial 28pt — tagline
  Column labels:    Calibri 9pt bold
  Code:             Calibri 7pt
  Annotations:      Calibri 9pt bold (red)
  Score tables:     Calibri 8pt
  Captions:         Calibri 10pt

Output: /home/baebs/thesis/review_mtgs/conflict_1-6_examples.pptx
"""

import json
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


REPO_ROOT     = Path(__file__).resolve().parent.parent
RESULTS_DIR   = REPO_ROOT / "scripts" / "results"
CONGRA_ROOT   = Path("/home/baebs/thesis/ConGra")
TEMPLATE_PATH = Path("/home/baebs/thesis/review_mtgs/conflict_1-6_examples.pptx")
OUT_PATH      = TEMPLATE_PATH

PILOTS = {
    ("qwen3",   "v1"):   "pilot_results_qwen3_v2.jsonl",
    ("qwen3",   "v2"):   "pilot_results_qwen3_skill-v2.jsonl",
    ("qwen3",   "v2.1"): "pilot_results_qwen3_skill-v2.1.jsonl",
    ("apertus", "v1"):   "pilot_results_apertus_v2.jsonl",
    ("apertus", "v2"):   "pilot_results_apertus_skill-v2.jsonl",
    ("apertus", "v2.1"): "pilot_results_apertus_skill-v2.1.jsonl",
}

RED_TITLE = RGBColor(0xD0, 0x21, 0x2C)
RED_ANNOT = RGBColor(0xB2, 0x18, 0x2B)
BLACK     = RGBColor(0, 0, 0)

# 4-col layout (Apertus + Qwen3 side-by-side)
COL_NORMAL = {
    "conflict": {"L": 0.20, "W": 1.80},
    "apertus":  {"L": 2.10, "W": 2.55},
    "qwen3":    {"L": 4.75, "W": 2.55},
    "gt":       {"L": 7.40, "W": 2.40},
}
# 5-col layout (one model, v1/v2/v2.1 comparison)
COL_SPLIT = {
    "conflict": {"L": 0.20, "W": 1.80},
    "v1":       {"L": 2.10, "W": 1.85},
    "v2":       {"L": 4.05, "W": 1.85},
    "v2.1":     {"L": 6.00, "W": 1.85},
    "gt":       {"L": 7.95, "W": 1.85},
}
ROW = {
    "title":      {"T": 0.05, "H": 0.55},
    "subtitle":   {"T": 0.60, "H": 0.85},
    "label":      {"T": 1.50, "H": 0.22},
    "code":       {"T": 1.75, "H": 2.20},   # ends 3.95
    "annotation": {"T": 4.00, "H": 0.25},
    "table":      {"T": 4.30, "H": 0.60},   # ends 4.90
    "caption":    {"T": 4.95, "H": 0.65},   # ends 5.60
}

# CASE specs — `split=True` produces 2 slides (one per model).
CASES = [
    {
        "case_id": "0xe63ff0ddae988357", "conflict_idx": 1,
        "tagline": "Apertus changes its pick under v2 — TF API → Keras API",
        "split": True,
        "per_model": {
            "apertus": {
                "annotations": {
                    "v1":   "tf.placeholder (wrong)",
                    "v2":   "K.placeholder (correct)",
                    "v2.1": "K.placeholder (correct)",
                },
                "caption": "v1 → v2 routes Apertus from TF API to Keras API "
                           "— the one clean pattern-routing change in the corpus.",
            },
            "qwen3": {
                "annotations": {
                    "v1":   "Drops from 0.836 to 0.437",
                    "v2":   "Recovers to 0.836",
                    "v2.1": "Regresses to 0.702",
                },
                "caption": "Qwen3 already at 0.836 without skill; v1 drops it, "
                           "v2 recovers, v2.1 regresses. Skill is "
                           "neutral-to-harmful for Qwen3 here.",
            },
        },
    },
    {
        "case_id": "0x96d20e6c", "conflict_idx": 1,
        "tagline": "v2 suppresses Apertus's over-generation",
        "split": False,
        "annotations": {
            "apertus": "v2 trims Apertus's over-generation",
            "qwen3":   "",
        },
        "caption_apertus": "v2 suppresses Apertus's over-generation. "
                           "Pattern routing stays wrong; metric improves "
                           "from v1 to v2.",
        "caption_qwen3":   "No skill helps Qwen3 — surrounding code would be "
                           "needed to solve this correctly; Qwen3 is good "
                           "in what it sees.",
    },
    {
        "case_id": "0xe4ff79aa", "conflict_idx": 1,
        "tagline": "Identifier-divergence (Pick) — metric vs. correctness diverge",
        "split": False,
        "annotations": {
            "apertus": "Picks pool_size correctly",
            "qwen3":   "Picks poolsize (wrong identifier)",
        },
        "caption_full": "Apertus picks pool_size (correct); Qwen3 picks "
                        "poolsize (wrong) — but the metric ranks Qwen3 higher "
                        "because its output is shorter (length-ratio dominates "
                        "the denominator).",
    },
    {
        "case_id": "0x8e6579cb86af64a8", "conflict_idx": 1,
        "tagline": "Same framing — opposite effects across models",
        "split": True,
        "per_model": {
            "apertus": {
                "annotations": {
                    "v1":   "0.375",
                    "v2":   "0.375",
                    "v2.1": "+0.21 → 0.584",
                },
                "caption": "v2.1's stronger output-discipline framing lifts "
                           "Apertus from 0.375 (v2) to 0.584 (v2.1).",
            },
            "qwen3": {
                "annotations": {
                    "v1":   "0.466",
                    "v2":   "0.672",
                    "v2.1": "−0.21 → 0.466",
                },
                "caption": "Same framing pushes Qwen3 to a shorter do-nothing "
                           "output; v2.1 drops Qwen3 from 0.672 back to 0.466.",
            },
        },
    },
    {
        "case_id": "0x32d8c89b39c2860b", "conflict_idx": 1,
        "tagline": "Apertus emits the same code under all 9 conditions",
        "split": False,
        "annotations": {
            "apertus": "Identical output every time",
            "qwen3":   "Wobbles to 0.512 under v2",
        },
        "caption_full": "Apertus produces the same output under all 9 "
                        "conditions (3 pilots × no-skill / sys / user). The "
                        "skill is a true no-op for Apertus. Qwen3 has zero "
                        "such cases — its outputs vary even when scores are tied.",
    },
    {
        "case_id": "0xd9272c5e0e8f15ee", "conflict_idx": 1,
        "tagline": "Task ceiling — and skill actively misleads Apertus",
        "split": False,
        "annotations": {
            "apertus": "Skill loses 0.11 vs no-skill",
            "qwen3":   "Stuck at 0.19 across all versions",
        },
        "caption_full": "Ground truth requires file-level pattern synthesis "
                        "outside the conflict region. No single-conflict "
                        "skill can recover this; on Apertus the skill actively "
                        "misleads (no-skill 0.289 → skill 0.177).",
    },
]


# ── ConGra loaders ─────────────────────────────────────────────────────────────
def load_conflict_and_answer(source_path, file_path, k, n=5):
    region_path = source_path.replace("merged_without_base", "regions") + ".region"
    regions = []
    for line in open(region_path):
        if "#" in line: continue
        s = line.strip()
        if s: regions.append(eval(s))
    region = regions[k - 1]
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        lines = f.read().split("\n")
    start, end = region[0] - 1, region[1]
    conflict_text = "\n".join(lines[start:end])
    resolved_path = source_path.replace("merged_without_base", "resolved").replace(
        "regions", "resolved")
    with open(resolved_path, "r", encoding="utf-8", errors="ignore") as f:
        rlines = f.read().split("\n")
    rstart, rend = region[2] - 1, region[3]
    resolved_text = "\n".join(rlines[max(0, rstart):rend])
    return conflict_text, resolved_text


def build_case_index():
    lang_root = CONGRA_ROOT / "data" / "congra_tiny_datasets" / "python"
    idx = {}
    for bucket_dir in sorted(lang_root.iterdir()):
        if not bucket_dir.is_dir(): continue
        meta = bucket_dir / "meta_list.txt"
        if not meta.is_file(): continue
        for line in meta.read_text().splitlines():
            line = line.strip()
            if not line or line.count(": ") != 2: continue
            sp, h, ci = line.split(": ")
            idx[h] = (bucket_dir.name, sp, int(ci))
    return idx


def normalize_case_id(short, index):
    if short in index: return short
    matches = [k for k in index if k.startswith(short)]
    if len(matches) == 1: return matches[0]
    raise KeyError(f"Cannot resolve {short!r}")


def load_pilots():
    out = {}
    for key, fname in PILOTS.items():
        by_case = defaultdict(dict)
        for line in open(RESULTS_DIR / fname):
            r = json.loads(line)
            by_case[(r["case_id"], r["conflict_idx"])][r["condition"]] = r
        out[key] = by_case
    return out


def fmt_score(x): return "—" if x is None else f"{x:.3f}"


def score_for(pilots, model, version, case_key):
    cond = "no-skill" if version == "no-skill" else f"skill-{version}-sys"
    pilot_key = "v2" if version == "no-skill" else version
    rec = pilots[(model, pilot_key)].get(case_key, {}).get(cond)
    if not rec or rec.get("error"): return None, None
    m = rec.get("metrics", {})
    return m.get("edit"), m.get("winnowing")


def get_solution(pilots, model, version, case_key):
    rec = pilots[(model, version)].get(case_key, {}).get(f"skill-{version}-sys")
    return (rec.get("resolution") or "") if rec and not rec.get("error") else ""


# ── Rich textbox + run-level highlight ────────────────────────────────────────
def set_run_highlight(run, hex_color):
    rPr = run._r.get_or_add_rPr()
    for el in rPr.findall(qn("a:highlight")):
        rPr.remove(el)
    hl = etree.SubElement(rPr, qn("a:highlight"))
    srgb = etree.SubElement(hl, qn("a:srgbClr"))
    srgb.set("val", hex_color)


def add_textbox_rich(slide, L, T, W, H, paragraphs, *,
                     font_name="Calibri", default_pt=10,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(L), Inches(T), Inches(W), Inches(H))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left   = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    first = True
    for paragraph in paragraphs:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        for r_spec in paragraph:
            run = p.add_run()
            run.text = r_spec.get("text", "")
            run.font.name = r_spec.get("font_name", font_name)
            run.font.size = Pt(r_spec.get("font_pt", default_pt))
            if r_spec.get("bold"):   run.font.bold = True
            if r_spec.get("italic"): run.font.italic = True
            color = r_spec.get("color")
            if color is not None:
                run.font.color.rgb = color
            hl = r_spec.get("highlight")
            if hl:
                set_run_highlight(run, hl)
    return box


def lines_to_paragraphs(text, **run_attrs):
    if not text:
        return [[{"text": "", **run_attrs}]]
    return [[{"text": line, **run_attrs}] for line in text.split("\n")]


# ── Score table ───────────────────────────────────────────────────────────────
def add_score_table(slide, L, T, W, H, scores):
    cols, rows = 5, 3
    shape = slide.shapes.add_table(rows, cols, Inches(L), Inches(T),
                                   Inches(W), Inches(H))
    tbl = shape.table
    headers = ["Scores", "no-skill", "v1", "v2", "v2.1"]
    data = [
        ["Edit Similarity", *[fmt_score(scores[v][0]) for v in ("no-skill","v1","v2","v2.1")]],
        ["Winnowing",       *[fmt_score(scores[v][1]) for v in ("no-skill","v1","v2","v2.1")]],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(8)
                run.font.bold = True
    for r_idx, row in enumerate(data, start=1):
        for c_idx, v in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = v
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(8)
    return shape


# ── Slide cleanup ─────────────────────────────────────────────────────────────
def clear_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    pres_part = prs.part
    for sld_id in list(sldIdLst):
        rId = sld_id.get(qn("r:id"))
        sldIdLst.remove(sld_id)
        pres_part.rels.pop(rId)


def find_layout_by_name(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(f"Layout {name!r} not found.")


def remove_placeholders(slide):
    for shape in list(slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)


def add_title_and_subtitle(slide, title_text, subtitle_text):
    add_textbox_rich(
        slide, 0.20, ROW["title"]["T"], 9.60, ROW["title"]["H"],
        [[{"text": title_text, "font_name": "Arial", "font_pt": 28,
           "bold": True, "color": RED_TITLE}]],
    )
    add_textbox_rich(
        slide, 0.20, ROW["subtitle"]["T"], 9.60, ROW["subtitle"]["H"],
        [[{"text": subtitle_text, "font_name": "Arial", "font_pt": 28,
           "color": BLACK}]],
    )


def add_label(slide, L, W, text):
    add_textbox_rich(
        slide, L, ROW["label"]["T"], W, ROW["label"]["H"],
        [[{"text": text, "font_pt": 9, "bold": True}]],
    )


def add_code(slide, L, W, text):
    code_attrs = {"font_name": "Calibri", "font_pt": 7}
    add_textbox_rich(
        slide, L, ROW["code"]["T"], W, ROW["code"]["H"],
        lines_to_paragraphs(text, **code_attrs),
        default_pt=7,
    )


def add_annotation(slide, L, W, text):
    add_textbox_rich(
        slide, L, ROW["annotation"]["T"], W, ROW["annotation"]["H"],
        [[{"text": text, "font_pt": 9, "bold": True, "color": RED_ANNOT}]],
        align=PP_ALIGN.CENTER,
    )


def add_caption(slide, L, W, text):
    add_textbox_rich(
        slide, L, ROW["caption"]["T"], W, ROW["caption"]["H"],
        [[{"text": text, "font_pt": 10}]],
    )


# ── Slide builders ────────────────────────────────────────────────────────────
def build_normal_slide(prs, layout, idx, case, pilots):
    """4-col layout: Conflict | Apertus | Qwen3 | GT."""
    slide = prs.slides.add_slide(layout)
    remove_placeholders(slide)
    case_key = (case["case_id"], case["conflict_idx"])

    add_title_and_subtitle(slide, f"Conflict {idx}", case["tagline"])

    for key, label in [("conflict", "Merge Conflict"),
                       ("apertus",  "Apertus Solution"),
                       ("qwen3",    "Qwen3 Solution"),
                       ("gt",       "Ground Truth")]:
        add_label(slide, COL_NORMAL[key]["L"], COL_NORMAL[key]["W"], label)

    add_code(slide, COL_NORMAL["conflict"]["L"], COL_NORMAL["conflict"]["W"],
             case["conflict_text"])
    add_code(slide, COL_NORMAL["apertus"]["L"], COL_NORMAL["apertus"]["W"],
             get_solution(pilots, "apertus", "v2.1", case_key) or "(empty)")
    add_code(slide, COL_NORMAL["qwen3"]["L"], COL_NORMAL["qwen3"]["W"],
             get_solution(pilots, "qwen3", "v2.1", case_key) or "(empty)")
    add_code(slide, COL_NORMAL["gt"]["L"], COL_NORMAL["gt"]["W"],
             case["ground_truth"])

    for key in ("apertus", "qwen3"):
        ann = (case.get("annotations") or {}).get(key, "")
        if ann:
            add_annotation(slide, COL_NORMAL[key]["L"],
                           COL_NORMAL[key]["W"], ann)

    a_scores = {v: score_for(pilots, "apertus", v, case_key)
                for v in ("no-skill", "v1", "v2", "v2.1")}
    q_scores = {v: score_for(pilots, "qwen3", v, case_key)
                for v in ("no-skill", "v1", "v2", "v2.1")}
    add_score_table(slide, COL_NORMAL["apertus"]["L"], ROW["table"]["T"],
                    COL_NORMAL["apertus"]["W"], ROW["table"]["H"], a_scores)
    add_score_table(slide, COL_NORMAL["qwen3"]["L"], ROW["table"]["T"],
                    COL_NORMAL["qwen3"]["W"], ROW["table"]["H"], q_scores)

    if "caption_full" in case:
        add_caption(slide, 0.20, 9.60, case["caption_full"])
    else:
        add_caption(slide, COL_NORMAL["apertus"]["L"], COL_NORMAL["apertus"]["W"],
                    case.get("caption_apertus", ""))
        add_caption(slide, COL_NORMAL["qwen3"]["L"], COL_NORMAL["qwen3"]["W"],
                    case.get("caption_qwen3", ""))


def build_split_slide(prs, layout, idx, case, pilots, model):
    """5-col layout: Conflict | v1 | v2 | v2.1 | GT, for one model."""
    slide = prs.slides.add_slide(layout)
    remove_placeholders(slide)
    case_key = (case["case_id"], case["conflict_idx"])
    model_label = "Apertus" if model == "apertus" else "Qwen3"
    pm = case["per_model"][model]

    add_title_and_subtitle(slide, f"Conflict {idx} — {model_label}", case["tagline"])

    add_label(slide, COL_SPLIT["conflict"]["L"], COL_SPLIT["conflict"]["W"],
              "Merge Conflict")
    add_label(slide, COL_SPLIT["v1"]["L"], COL_SPLIT["v1"]["W"],
              f"{model_label} — v1")
    add_label(slide, COL_SPLIT["v2"]["L"], COL_SPLIT["v2"]["W"],
              f"{model_label} — v2")
    add_label(slide, COL_SPLIT["v2.1"]["L"], COL_SPLIT["v2.1"]["W"],
              f"{model_label} — v2.1")
    add_label(slide, COL_SPLIT["gt"]["L"], COL_SPLIT["gt"]["W"],
              "Ground Truth")

    add_code(slide, COL_SPLIT["conflict"]["L"], COL_SPLIT["conflict"]["W"],
             case["conflict_text"])
    for v in ("v1", "v2", "v2.1"):
        text = get_solution(pilots, model, v, case_key) or "(empty)"
        add_code(slide, COL_SPLIT[v]["L"], COL_SPLIT[v]["W"], text)
    add_code(slide, COL_SPLIT["gt"]["L"], COL_SPLIT["gt"]["W"],
             case["ground_truth"])

    for v in ("v1", "v2", "v2.1"):
        ann = pm.get("annotations", {}).get(v, "")
        if ann:
            add_annotation(slide, COL_SPLIT[v]["L"], COL_SPLIT[v]["W"], ann)

    # One score table for this model, spanning v1..v2.1 columns
    scores = {v: score_for(pilots, model, v, case_key)
              for v in ("no-skill", "v1", "v2", "v2.1")}
    table_L = COL_SPLIT["v1"]["L"]
    table_W = (COL_SPLIT["v2.1"]["L"] + COL_SPLIT["v2.1"]["W"]) - COL_SPLIT["v1"]["L"]
    add_score_table(slide, table_L, ROW["table"]["T"],
                    table_W, ROW["table"]["H"], scores)

    add_caption(slide, 0.20, 9.60, pm.get("caption", ""))


# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    print("Loading case index...")
    case_idx = build_case_index()
    print(f"  {len(case_idx)} python cases indexed")

    print("Loading pilot JSONLs...")
    pilots = load_pilots()

    print("\nResolving cases:")
    resolved = []
    for spec in CASES:
        cid = normalize_case_id(spec["case_id"], case_idx)
        bucket, src_rel, _ = case_idx[cid]
        source_path = str(CONGRA_ROOT / "data" / "raw_datasets" / src_rel)
        hash_file = str(CONGRA_ROOT / "data" / "congra_tiny_datasets"
                        / "python" / bucket / cid)
        conflict_text, gt = load_conflict_and_answer(
            source_path, hash_file, spec["conflict_idx"], n=5)
        resolved.append({**spec, "case_id": cid, "bucket": bucket,
                         "conflict_text": conflict_text, "ground_truth": gt})
        print(f"  {cid[:18]}... bucket={bucket:>16} "
              f"conflict={len(conflict_text.splitlines())}L "
              f"GT={len(gt.splitlines())}L")

    prs = Presentation(str(TEMPLATE_PATH))
    print("Clearing existing slides...")
    clear_all_slides(prs)

    layout = find_layout_by_name(prs, "3: Titel-Folie ohne Bild")
    print(f"Using layout: {layout.name}")

    n_slides = 0
    for i, case in enumerate(resolved, start=1):
        if case.get("split"):
            build_split_slide(prs, layout, i, case, pilots, "apertus")
            n_slides += 1
            print(f"  slide {n_slides}: Conflict {i} — Apertus built")
            build_split_slide(prs, layout, i, case, pilots, "qwen3")
            n_slides += 1
            print(f"  slide {n_slides}: Conflict {i} — Qwen3 built")
        else:
            build_normal_slide(prs, layout, i, case, pilots)
            n_slides += 1
            print(f"  slide {n_slides}: Conflict {i} built")

    print(f"\nSaving: {OUT_PATH} ({n_slides} slides)")
    prs.save(str(OUT_PATH))
    print("Done.")


if __name__ == "__main__":
    main()
