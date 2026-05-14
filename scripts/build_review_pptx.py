"""
build_review_pptx.py — assemble a Uni Bern review-meeting pptx (issue #48).

Default args reproduce the 2026-05-22 deck:
  - Carry the first 7 slides verbatim from Thesis_Review_2026-05-08.pptx.
  - Drop slides 8..14 (slide 8 text-table is being replaced; 9..14 are being
    rebuilt as part of a separate body-deck redesign).
  - Append an Appendix title slide (layout "3: Titel-Folie ohne Bild").
  - Append a per-case-outcomes chart slide (layout "7a: Folie fur Grafik/Bild")
    embedding docs/figures/per_case_outcomes_v1_v2_v21.png.

Reusable for future reviews: override --source/--output/--keep-slides and the
text knobs to assemble a different deck.

Layout placeholder indices used (from reference_pptx_template / inspected in
the source deck):
  Layout 3: idx=0 Titel A (red), idx=11 Titel B (black), idx=12 Date+location.
  Layout 7a: idx=20 TOC header, idx=0 Titel A, idx=11 Titel B, idx=1 chart area.

python-pptx slide-deletion gotcha: removing from sldIdLst alone leaves orphan
parts; we drop the relationship too (see delete_slides_after).
"""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn


REPO_ROOT     = Path(__file__).resolve().parent.parent
DEFAULT_SRC   = Path("/home/baebs/thesis/review_mtgs/Thesis_Review_2026-05-08.pptx")
DEFAULT_OUT   = Path("/home/baebs/thesis/review_mtgs/Thesis_Review_2026-05-22.pptx")
DEFAULT_CHART = REPO_ROOT / "docs" / "figures" / "per_case_outcomes_v1_v2_v21.png"


def find_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(f"Layout {name!r} not found in any slide master.")


def set_ph_text(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text_frame.text = text
            return
    raise KeyError(f"No placeholder with idx={idx} on slide.")


def drop_placeholder(slide, idx):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == idx:
            sp = ph._element
            sp.getparent().remove(sp)
            return


def delete_slides_after(prs, keep_n):
    sldIdLst = prs.slides._sldIdLst
    pres_part = prs.part
    for sld_id in list(sldIdLst)[keep_n:]:
        rId = sld_id.get(qn("r:id"))
        sldIdLst.remove(sld_id)
        pres_part.rels.pop(rId)


def add_appendix_title(prs, title_a, title_b, date_loc):
    layout = find_layout(prs, "3: Titel-Folie ohne Bild")
    slide = prs.slides.add_slide(layout)
    set_ph_text(slide, 0, title_a)
    set_ph_text(slide, 11, title_b)
    set_ph_text(slide, 12, date_loc)
    drop_placeholder(slide, 10)   # moderator line — leave clean
    return slide


def add_chart_slide(prs, title_a, title_b, toc_header, image_path):
    layout = find_layout(prs, "7a: Folie für Grafik/Bild")
    slide = prs.slides.add_slide(layout)
    set_ph_text(slide, 0, title_a)
    set_ph_text(slide, 11, title_b)
    set_ph_text(slide, 20, toc_header)

    drop_placeholder(slide, 1)   # remove "Media Platzhalter" hint text

    # Chart area on layout 7a: L=0, T=2.11", W=10", H=3.35".
    # Add by height only so aspect is preserved, then center horizontally.
    pic = slide.shapes.add_picture(
        str(image_path),
        left=Inches(0),
        top=Inches(2.11),
        height=Inches(3.35),
    )
    pic.left = (prs.slide_width - pic.width) // 2
    return slide


def main():
    p = argparse.ArgumentParser(
        description="Assemble Thesis_Review_YYYY-MM-DD.pptx from a prior deck.",
    )
    p.add_argument("--source", type=Path, default=DEFAULT_SRC,
                   help="Source pptx to carry leading slides from")
    p.add_argument("--output", type=Path, default=DEFAULT_OUT,
                   help="Output pptx path")
    p.add_argument("--keep-slides", type=int, default=7,
                   help="Number of leading slides to carry from source (default 7)")
    p.add_argument("--appendix-title-a", default="Appendix")
    p.add_argument("--appendix-title-b", default="Per-case outcome details")
    p.add_argument("--appendix-date", default="2026-05-22, Bern")
    p.add_argument("--chart-title-a", default="Per-case outcome distribution")
    p.add_argument("--chart-title-b", default="n=20, python/func")
    p.add_argument("--chart-toc-header", default="Appendix")
    p.add_argument("--chart-image", type=Path, default=DEFAULT_CHART)
    args = p.parse_args()

    if not args.source.is_file():
        p.error(f"Source pptx not found: {args.source}")
    if not args.chart_image.is_file():
        p.error(f"Chart image not found: {args.chart_image}")

    prs = Presentation(str(args.source))
    n_orig = len(prs.slides)
    print(f"Loaded source: {args.source} ({n_orig} slides)")

    delete_slides_after(prs, args.keep_slides)
    print(f"  -> kept slides 1..{args.keep_slides}; dropped {n_orig - args.keep_slides}")

    add_appendix_title(
        prs, args.appendix_title_a, args.appendix_title_b, args.appendix_date,
    )
    print("  -> appended appendix title slide (layout 3)")

    add_chart_slide(
        prs, args.chart_title_a, args.chart_title_b,
        args.chart_toc_header, args.chart_image,
    )
    print(f"  -> appended chart slide (layout 7a) with {args.chart_image.name}")

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output))
    print(f"Wrote: {args.output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
